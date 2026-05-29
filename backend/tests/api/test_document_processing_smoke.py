from collections.abc import Generator
from contextlib import suppress
from datetime import UTC, datetime
from io import BytesIO

import pytest
from fastapi.testclient import TestClient
from sqlalchemy.pool import StaticPool
from sqlmodel import Session, SQLModel, create_engine

from app.core.config import get_settings
from app.db.session import get_session
from app.main import app
from app.models.user import User
from app.services.users import DEFAULT_USER_ID

# 本文件用 API 烟测覆盖上传、文档处理、分块读取和失败文档列表展示的完整链路。


@pytest.fixture
def session() -> Generator[Session]:
    with suppress(ModuleNotFoundError):
        import app.models.document  # noqa: F401

    engine = create_engine(
        "sqlite://",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    SQLModel.metadata.create_all(engine)

    with Session(engine) as test_session:
        yield test_session


@pytest.fixture
def smoke_client(monkeypatch, session: Session, tmp_path) -> Generator[TestClient]:
    monkeypatch.setenv("UPLOAD_STORAGE_DIR", str(tmp_path))
    monkeypatch.setenv(
        "ALLOWED_UPLOAD_CONTENT_TYPES",
        ",".join(
            [
                "text/plain",
                "text/markdown",
                "application/pdf",
                "application/vnd.ms-powerpoint",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ]
        ),
    )
    get_settings.cache_clear()

    def override_get_session() -> Generator[Session]:
        yield session

    app.dependency_overrides[get_session] = override_get_session

    with TestClient(app) as test_client:
        yield test_client

    app.dependency_overrides.clear()
    get_settings.cache_clear()


def seed_current_user(session: Session) -> None:
    created_at = datetime(2026, 5, 28, tzinfo=UTC)
    session.add(
        User(
            id=DEFAULT_USER_ID,
            display_name="Default User",
            status="active",
            created_at=created_at,
            updated_at=created_at,
        )
    )
    session.commit()


def make_pdf_bytes(text: str = "PDF smoke note about citations") -> bytes:
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R "
            b"/Resources << /Font << /F1 4 0 R >> >> "
            b"/MediaBox [0 0 612 792] /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(stream), stream),
    ]
    return build_pdf(objects)


def make_blank_pdf_bytes() -> bytes:
    return build_pdf(
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            (b"<< /Type /Page /Parent 2 0 R /Resources << >> /MediaBox [0 0 612 792] >>"),
        ]
    )


def build_pdf(objects: list[bytes]) -> bytes:
    pdf = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{index} 0 obj\n".encode("ascii"))
        pdf.extend(obj)
        pdf.extend(b"\nendobj\n")

    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode("ascii"))
    for offset in offsets:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    pdf.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    return bytes(pdf)


def make_docx_bytes() -> bytes:
    from docx import Document as DocxDocument

    document = DocxDocument()
    document.add_heading("DOCX smoke note", level=1)
    document.add_paragraph("DOCX body about chunk creation.")
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def make_pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PowerPoint smoke note"
    slide.placeholders[1].text = "Slide body about traceable chunks."
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def upload_and_process(
    client: TestClient,
    *,
    filename: str,
    content: bytes,
    content_type: str,
) -> dict:
    upload_response = client.post(
        "/api/uploads",
        files={"file": (filename, content, content_type)},
    )
    assert upload_response.status_code == 201

    uploaded_file_id = upload_response.json()["id"]
    document_response = client.post(
        "/api/documents",
        json={"uploaded_file_id": uploaded_file_id},
    )
    assert document_response.status_code == 201

    duplicate_response = client.post(
        "/api/documents",
        json={"uploaded_file_id": uploaded_file_id},
    )
    assert duplicate_response.status_code == 409
    assert duplicate_response.json()["existing_document"]["uploaded_file_id"] == uploaded_file_id

    return document_response.json()


# 测试首批支持格式可从上传一路处理成文档分块，同时扫描件 PDF 会以 failed 状态保留错误原因。
def test_document_processing_api_smoke_for_supported_formats_and_failed_pdf(
    smoke_client: TestClient,
    session: Session,
) -> None:
    seed_current_user(session)
    pptx_content = make_pptx_bytes()
    samples = [
        ("notes.txt", b"TXT smoke note about retrieval", "text/plain"),
        ("notes.md", b"# Markdown smoke note\n\nChunk this document.", "text/markdown"),
        ("notes.pdf", make_pdf_bytes(), "application/pdf"),
        (
            "notes.docx",
            make_docx_bytes(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
        (
            "notes.pptx",
            pptx_content,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        ("notes.ppt", pptx_content, "application/vnd.ms-powerpoint"),
    ]

    parsed_documents = [
        upload_and_process(
            smoke_client,
            filename=filename,
            content=content,
            content_type=content_type,
        )
        for filename, content, content_type in samples
    ]

    for document in parsed_documents:
        assert document["status"] == "parsed"
        assert document["chunk_count"] > 0
        chunks_response = smoke_client.get(f"/api/documents/{document['id']}/chunks")
        assert chunks_response.status_code == 200
        assert len(chunks_response.json()) == document["chunk_count"]

    failed_document = upload_and_process(
        smoke_client,
        filename="scan.pdf",
        content=make_blank_pdf_bytes(),
        content_type="application/pdf",
    )
    assert failed_document["status"] == "failed"
    assert failed_document["chunk_count"] == 0
    assert "OCR" in failed_document["error_message"]

    documents_response = smoke_client.get("/api/documents")
    assert documents_response.status_code == 200
    statuses = {document["status"] for document in documents_response.json()}
    assert {"parsed", "failed"} <= statuses

    failed_chunks_response = smoke_client.get(f"/api/documents/{failed_document['id']}/chunks")
    assert failed_chunks_response.status_code == 200
    assert failed_chunks_response.json() == []

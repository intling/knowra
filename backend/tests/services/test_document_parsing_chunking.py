from collections.abc import Callable
from hashlib import sha256
from importlib import import_module
from io import BytesIO

import pytest

# 本文件验证默认解析器、BPE tokenizer 和分块器对首批文档格式的基础处理能力。


def get_processing_module():
    return import_module("app.services.document_processing")


def make_pdf_bytes(text: str = "PDF lecture note about retrieval") -> bytes:
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R "
            b"/Resources << /Font << /F1 4 0 R >> >> "
            b"/MediaBox [0 0 612 792] /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(stream), stream),
    ]
    pdf = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{index} 0 obj\n".encode("ascii"))
        pdf.extend(obj)
        pdf.extend(b"\nendobj\n")

    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode("ascii"))
    for offset in offsets:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    pdf.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    return bytes(pdf)


def make_docx_bytes() -> bytes:
    from docx import Document as DocxDocument

    document = DocxDocument()
    document.add_heading("DOCX lecture note", level=1)
    document.add_paragraph("DOCX body paragraph about semantic indexing.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Term"
    table.cell(0, 1).text = "Definition"

    output = BytesIO()
    document.save(output)
    return output.getvalue()


def make_pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPTX lecture note"
    slide.placeholders[1].text = "Slide body about chunk provenance."
    slide.notes_slide.notes_text_frame.text = "Speaker note about citations."

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


SAMPLES: list[tuple[str, str, Callable[[], bytes], str, str, str]] = [
    (
        "course-notes.txt",
        "text/plain",
        lambda: b"TXT lecture note about private knowledge",
        "TXT lecture note",
        "line_start",
        "source_locator_json",
    ),
    (
        "course-notes.md",
        "text/markdown",
        lambda: b"# Markdown lecture note\n\nDetails about retrieval.",
        "Markdown lecture note",
        "heading_path",
        "metadata_json",
    ),
    (
        "course-notes.pdf",
        "application/pdf",
        make_pdf_bytes,
        "PDF lecture note",
        "page",
        "source_locator_json",
    ),
    (
        "course-notes.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        make_docx_bytes,
        "DOCX lecture note",
        "paragraph",
        "source_locator_json",
    ),
    (
        "course-notes.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        make_pptx_bytes,
        "PPTX lecture note",
        "slide",
        "source_locator_json",
    ),
]


# 测试 TXT/Markdown/PDF/DOCX/PPTX 都能被解析并生成带来源定位和哈希的分块。
@pytest.mark.parametrize(
    ("filename", "content_type", "content_factory", "expected_text", "locator_key", "field_name"),
    SAMPLES,
)
def test_default_parsers_create_chunks_for_first_batch_formats(
    filename: str,
    content_type: str,
    content_factory: Callable[[], bytes],
    expected_text: str,
    locator_key: str,
    field_name: str,
) -> None:
    processing = get_processing_module()
    registry = processing.create_default_parser_registry()
    tokenizer = processing.BpeTokenizer()
    chunker = processing.BpeChunker(tokenizer=tokenizer, max_tokens=32, overlap_tokens=4)

    parsed = registry.parse(
        filename=filename,
        content_type=content_type,
        content=content_factory(),
    )
    chunks = chunker.chunk(parsed)

    assert expected_text in parsed.text
    assert parsed.parser_name
    assert parsed.parser_version
    assert chunks
    assert [chunk.chunk_index for chunk in chunks] == list(range(len(chunks)))

    first_chunk = chunks[0]
    assert expected_text in first_chunk.content
    assert first_chunk.content_sha256 == sha256(first_chunk.content.encode()).hexdigest()
    assert first_chunk.char_start == 0
    assert first_chunk.char_end > first_chunk.char_start
    assert 0 < first_chunk.token_count <= 32
    assert first_chunk.source_locator_json
    assert first_chunk.metadata_json is not None

    field_value = getattr(first_chunk, field_name)
    assert locator_key in field_value


# 测试 BPE tokenizer 暴露名称版本，并在相同输入下生成稳定的重叠窗口。
def test_bpe_tokenizer_reports_name_version_and_deterministic_windows() -> None:
    processing = get_processing_module()
    tokenizer = processing.BpeTokenizer()

    text = "alpha beta gamma delta epsilon zeta eta theta"
    first_windows = tokenizer.windows(text, max_tokens=4, overlap_tokens=1)
    second_windows = tokenizer.windows(text, max_tokens=4, overlap_tokens=1)

    assert tokenizer.name
    assert tokenizer.version
    assert tokenizer.count(text) > 4
    assert first_windows == second_windows
    assert all(window.token_count <= 4 for window in first_windows)
    assert first_windows[0].char_start == 0

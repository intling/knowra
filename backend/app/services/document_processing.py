from dataclasses import dataclass
from hashlib import sha256
from importlib.metadata import PackageNotFoundError, version
from io import BytesIO
from pathlib import Path
from typing import Any, Protocol
from uuid import UUID

import tiktoken
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader
from sqlmodel import Session, select

from app.models.document import Document, DocumentChunk
from app.models.uploaded_file import UploadedFile
from app.models.user import User
from app.services.uploads import LocalFileStorage, UploadStorageError


class DocumentProcessingServiceError(Exception):
    pass


class UploadedFileNotFoundError(DocumentProcessingServiceError):
    pass


class UploadedFileNotStoredError(DocumentProcessingServiceError):
    pass


class DocumentAlreadyExistsError(DocumentProcessingServiceError):
    def __init__(self, existing_document: Document) -> None:
        super().__init__("Document already exists")
        self.existing_document = existing_document


class UnsupportedDocumentForProcessingError(DocumentProcessingServiceError):
    def __init__(self, message: str, failed_document: Document) -> None:
        super().__init__(message)
        self.failed_document = failed_document


class DocumentParseError(Exception):
    pass


class UnsupportedDocumentTypeError(DocumentParseError):
    pass


@dataclass(frozen=True)
class ParsedSegment:
    text: str
    char_start: int
    char_end: int
    source_locator_json: dict[str, Any]
    metadata_json: dict[str, Any]


@dataclass(frozen=True)
class ParsedDocument:
    text: str
    segments: list[ParsedSegment]
    metadata_json: dict[str, Any]
    parser_name: str
    parser_version: str


@dataclass(frozen=True)
class TokenWindow:
    content: str
    token_count: int
    char_start: int
    char_end: int


@dataclass(frozen=True)
class DocumentChunkDraft:
    chunk_index: int
    content: str
    content_sha256: str
    char_start: int
    char_end: int
    token_count: int
    source_locator_json: dict[str, Any]
    metadata_json: dict[str, Any]


class Parser(Protocol):
    name: str
    version: str
    supported_content_types: set[str]
    supported_extensions: set[str]

    def parse(
        self, *, filename: str, content_type: str | None, content: bytes
    ) -> ParsedDocument: ...


class ParserRegistry:
    def __init__(self, parsers: list[Parser]) -> None:
        self.parsers = parsers

    def parse(self, *, filename: str, content_type: str | None, content: bytes) -> ParsedDocument:
        parser = self.get_parser(filename=filename, content_type=content_type)
        if parser is None:
            raise UnsupportedDocumentTypeError("Unsupported document type")

        try:
            return parser.parse(filename=filename, content_type=content_type, content=content)
        except DocumentParseError:
            raise
        except Exception as exc:
            raise DocumentParseError("Failed to parse document") from exc

    def get_parser(self, *, filename: str, content_type: str | None) -> Parser | None:
        extension = Path(filename).suffix.lower()
        normalized_content_type = (content_type or "").lower()
        for parser in self.parsers:
            if (
                normalized_content_type
                and normalized_content_type in parser.supported_content_types
            ):
                return parser
            if extension and extension in parser.supported_extensions:
                return parser

        return None


class TextParser:
    name = "text"
    version = "1"
    supported_content_types = {"text/plain"}
    supported_extensions = {".txt"}

    def parse(self, *, filename: str, content_type: str | None, content: bytes) -> ParsedDocument:
        text = decode_utf8(content)
        segments = build_line_segments(text)
        return ParsedDocument(
            text=text,
            segments=segments,
            metadata_json={"filename": filename},
            parser_name=self.name,
            parser_version=self.version,
        )


class MarkdownParser:
    name = "markdown"
    version = "1"
    supported_content_types = {"text/markdown", "text/x-markdown"}
    supported_extensions = {".md", ".markdown"}

    def parse(self, *, filename: str, content_type: str | None, content: bytes) -> ParsedDocument:
        text = decode_utf8(content)
        heading_path = extract_markdown_heading_path(text)
        segment = ParsedSegment(
            text=text,
            char_start=0,
            char_end=len(text),
            source_locator_json={"line_start": 1},
            metadata_json={"heading_path": heading_path},
        )
        return ParsedDocument(
            text=text,
            segments=[segment],
            metadata_json={"filename": filename, "heading_path": heading_path},
            parser_name=self.name,
            parser_version=self.version,
        )


class PdfParser:
    name = "pdf"
    version = "1"
    supported_content_types = {"application/pdf"}
    supported_extensions = {".pdf"}

    def parse(self, *, filename: str, content_type: str | None, content: bytes) -> ParsedDocument:
        try:
            reader = PdfReader(BytesIO(content))
        except Exception as exc:
            raise DocumentParseError("Failed to parse PDF") from exc

        page_items: list[tuple[str, dict[str, Any], dict[str, Any]]] = []
        for page_index, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if page_text:
                page_items.append((page_text, {"page": page_index}, {}))

        if not page_items:
            raise DocumentParseError("PDF has no extractable text and may require OCR")

        text, segments = build_structured_text(page_items)
        return ParsedDocument(
            text=text,
            segments=segments,
            metadata_json={"filename": filename, "page_count": len(reader.pages)},
            parser_name=self.name,
            parser_version=self.version,
        )


class DocxParser:
    name = "docx"
    version = "1"
    supported_content_types = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }
    supported_extensions = {".docx"}

    def parse(self, *, filename: str, content_type: str | None, content: bytes) -> ParsedDocument:
        try:
            document = DocxDocument(BytesIO(content))
        except Exception as exc:
            raise DocumentParseError("Failed to parse DOCX") from exc

        items: list[tuple[str, dict[str, Any], dict[str, Any]]] = []
        heading_path: list[str] = []
        for index, paragraph in enumerate(document.paragraphs):
            paragraph_text = paragraph.text.strip()
            if not paragraph_text:
                continue

            style_name = paragraph.style.name if paragraph.style is not None else ""
            if style_name.lower().startswith("heading"):
                heading_path = [paragraph_text]

            items.append(
                (
                    paragraph_text,
                    {"paragraph": index},
                    {"heading_path": heading_path.copy()},
                )
            )

        for table_index, table in enumerate(document.tables):
            rows: list[str] = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    rows.append(row_text)
            if rows:
                items.append(
                    (
                        "\n".join(rows),
                        {"structure_path": f"table[{table_index}]"},
                        {"heading_path": heading_path.copy()},
                    )
                )

        if not items:
            raise DocumentParseError("DOCX parsed text is empty")

        text, segments = build_structured_text(items)
        return ParsedDocument(
            text=text,
            segments=segments,
            metadata_json={"filename": filename},
            parser_name=self.name,
            parser_version=self.version,
        )


class PowerPointParser:
    name = "powerpoint"
    version = "1"
    supported_content_types = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }
    supported_extensions = {".pptx", ".ppt"}

    def parse(self, *, filename: str, content_type: str | None, content: bytes) -> ParsedDocument:
        try:
            presentation = Presentation(BytesIO(content))
        except Exception as exc:
            raise DocumentParseError("Failed to parse PowerPoint document") from exc

        items: list[tuple[str, dict[str, Any], dict[str, Any]]] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                shape_text = getattr(shape, "text", "").strip()
                if shape_text:
                    slide_texts.append(shape_text)

            notes_text = ""
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(notes_text)

            if slide_texts:
                items.append(("\n".join(slide_texts), {"slide": slide_index}, {}))

        if not items:
            raise DocumentParseError("PowerPoint parsed text is empty")

        text, segments = build_structured_text(items)
        return ParsedDocument(
            text=text,
            segments=segments,
            metadata_json={"filename": filename, "slide_count": len(presentation.slides)},
            parser_name=self.name,
            parser_version=self.version,
        )


class BpeTokenizer:
    def __init__(self, encoding_name: str = "cl100k_base") -> None:
        self.encoding_name = encoding_name
        self.encoding = tiktoken.get_encoding(encoding_name)
        self.name = f"tiktoken:{encoding_name}"
        self.version = package_version("tiktoken")

    def count(self, text: str) -> int:
        return len(self.encoding.encode(text))

    def windows(self, text: str, *, max_tokens: int, overlap_tokens: int = 0) -> list[TokenWindow]:
        if max_tokens <= 0:
            raise ValueError("max_tokens must be greater than zero")
        if overlap_tokens < 0 or overlap_tokens >= max_tokens:
            raise ValueError("overlap_tokens must be smaller than max_tokens")

        tokens = self.encoding.encode(text)
        if not tokens:
            return []

        windows: list[TokenWindow] = []
        step = max_tokens - overlap_tokens
        start = 0
        while start < len(tokens):
            end = min(start + max_tokens, len(tokens))
            window_tokens = tokens[start:end]
            window_text = self.encoding.decode(window_tokens)
            char_start = len(self.encoding.decode(tokens[:start]))
            char_end = len(self.encoding.decode(tokens[:end]))
            windows.append(
                TokenWindow(
                    content=window_text,
                    token_count=len(window_tokens),
                    char_start=char_start,
                    char_end=char_end,
                )
            )
            if end == len(tokens):
                break
            start += step

        return windows


class BpeChunker:
    name = "bpe-window"
    version = "1"

    def __init__(
        self,
        *,
        tokenizer: BpeTokenizer,
        max_tokens: int = 800,
        overlap_tokens: int = 80,
    ) -> None:
        self.tokenizer = tokenizer
        self.max_tokens = max_tokens
        self.overlap_tokens = overlap_tokens

    def chunk(self, parsed: ParsedDocument) -> list[DocumentChunkDraft]:
        if not parsed.text.strip():
            raise DocumentParseError("Parsed document text is empty")

        chunks: list[DocumentChunkDraft] = []
        for window in self.tokenizer.windows(
            parsed.text,
            max_tokens=self.max_tokens,
            overlap_tokens=self.overlap_tokens,
        ):
            if not window.content.strip():
                continue

            overlapping_segments = [
                segment
                for segment in parsed.segments
                if segment.char_start < window.char_end and segment.char_end > window.char_start
            ]
            source_locator_json = combine_source_locators(overlapping_segments)
            metadata_json = combine_metadata(parsed.metadata_json, overlapping_segments)
            chunks.append(
                DocumentChunkDraft(
                    chunk_index=len(chunks),
                    content=window.content,
                    content_sha256=sha256(window.content.encode()).hexdigest(),
                    char_start=window.char_start,
                    char_end=window.char_end,
                    token_count=window.token_count,
                    source_locator_json=source_locator_json,
                    metadata_json=metadata_json,
                )
            )

        if not chunks:
            raise DocumentParseError("Parsed document text is empty")

        return chunks


class DocumentProcessingService:
    def __init__(
        self,
        *,
        session: Session,
        storage: LocalFileStorage,
        parser_registry: ParserRegistry | None = None,
        tokenizer: BpeTokenizer | None = None,
        chunker: BpeChunker | None = None,
    ) -> None:
        self.session = session
        self.storage = storage
        self.parser_registry = parser_registry or create_default_parser_registry()
        self.tokenizer = tokenizer or BpeTokenizer()
        self.chunker = chunker or BpeChunker(tokenizer=self.tokenizer)

    def create_document(self, *, current_user: User, uploaded_file_id: UUID) -> Document:
        existing_document = self.session.exec(
            select(Document).where(Document.uploaded_file_id == uploaded_file_id)
        ).first()
        if existing_document is not None:
            raise DocumentAlreadyExistsError(existing_document)

        uploaded_file = self.session.exec(
            select(UploadedFile).where(
                UploadedFile.id == uploaded_file_id,
                UploadedFile.owner_user_id == current_user.id,
            )
        ).first()
        if uploaded_file is None:
            raise UploadedFileNotFoundError("Uploaded file not found")
        if uploaded_file.status != "stored":
            raise UploadedFileNotStoredError("Uploaded file must be stored before processing")

        try:
            with self.storage.open(uploaded_file.storage_key) as original_file:
                content = original_file.read()
            parsed = self.parser_registry.parse(
                filename=uploaded_file.original_filename,
                content_type=uploaded_file.content_type,
                content=content,
            )
            chunk_drafts = self.chunker.chunk(parsed)
        except UnsupportedDocumentTypeError as exc:
            failed_document = self.create_failed_document(
                current_user=current_user,
                uploaded_file=uploaded_file,
                error_message=str(exc),
            )
            raise UnsupportedDocumentForProcessingError(str(exc), failed_document) from exc
        except (DocumentParseError, UploadStorageError) as exc:
            return self.create_failed_document(
                current_user=current_user,
                uploaded_file=uploaded_file,
                error_message=str(exc),
            )

        document = Document(
            owner_user_id=current_user.id,
            uploaded_file_id=uploaded_file.id,
            title=uploaded_file.original_filename,
            source_content_type=uploaded_file.content_type,
            parser_name=parsed.parser_name,
            parser_version=parsed.parser_version,
            chunker_name=self.chunker.name,
            chunker_version=self.chunker.version,
            tokenizer_name=self.tokenizer.name,
            tokenizer_version=self.tokenizer.version,
            status="parsed",
            chunk_count=len(chunk_drafts),
            total_chars=len(parsed.text),
            content_sha256=sha256(parsed.text.encode()).hexdigest(),
            metadata_json=parsed.metadata_json,
            error_message=None,
        )
        self.session.add(document)

        try:
            self.session.flush()
            for draft in chunk_drafts:
                self.session.add(
                    DocumentChunk(
                        document_id=document.id,
                        owner_user_id=current_user.id,
                        chunk_index=draft.chunk_index,
                        content=draft.content,
                        content_sha256=draft.content_sha256,
                        char_start=draft.char_start,
                        char_end=draft.char_end,
                        token_count=draft.token_count,
                        source_locator_json=draft.source_locator_json,
                        metadata_json=draft.metadata_json,
                    )
                )
            self.session.commit()
            self.session.refresh(document)
        except Exception as exc:
            self.session.rollback()
            raise DocumentProcessingServiceError("Failed to save parsed document") from exc

        return document

    def create_failed_document(
        self,
        *,
        current_user: User,
        uploaded_file: UploadedFile,
        error_message: str,
    ) -> Document:
        document = Document(
            owner_user_id=current_user.id,
            uploaded_file_id=uploaded_file.id,
            title=uploaded_file.original_filename,
            source_content_type=uploaded_file.content_type,
            parser_name=None,
            parser_version=None,
            chunker_name=self.chunker.name,
            chunker_version=self.chunker.version,
            tokenizer_name=self.tokenizer.name,
            tokenizer_version=self.tokenizer.version,
            status="failed",
            chunk_count=0,
            total_chars=0,
            content_sha256=None,
            metadata_json={},
            error_message=error_message,
        )
        self.session.add(document)

        try:
            self.session.commit()
            self.session.refresh(document)
        except Exception as exc:
            self.session.rollback()
            raise DocumentProcessingServiceError("Failed to save failed document") from exc

        return document

    def list_documents(self, *, current_user: User) -> list[Document]:
        return list(
            self.session.exec(
                select(Document)
                .where(Document.owner_user_id == current_user.id)
                .order_by(Document.created_at)
            ).all()
        )

    def get_document(self, *, current_user: User, document_id: UUID) -> Document | None:
        return self.session.exec(
            select(Document).where(
                Document.id == document_id,
                Document.owner_user_id == current_user.id,
            )
        ).first()

    def list_chunks(self, *, current_user: User, document_id: UUID) -> list[DocumentChunk] | None:
        document = self.get_document(current_user=current_user, document_id=document_id)
        if document is None:
            return None
        if document.status != "parsed":
            return []

        return list(
            self.session.exec(
                select(DocumentChunk)
                .where(
                    DocumentChunk.document_id == document_id,
                    DocumentChunk.owner_user_id == current_user.id,
                )
                .order_by(DocumentChunk.chunk_index)
            ).all()
        )


def create_default_parser_registry() -> ParserRegistry:
    return ParserRegistry(
        [
            TextParser(),
            MarkdownParser(),
            PdfParser(),
            DocxParser(),
            PowerPointParser(),
        ]
    )


def decode_utf8(content: bytes) -> str:
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise DocumentParseError("TXT/Markdown content must be valid UTF-8") from exc


def build_line_segments(text: str) -> list[ParsedSegment]:
    segments: list[ParsedSegment] = []
    offset = 0
    for line_number, line in enumerate(text.splitlines(keepends=True), start=1):
        line_start = offset
        line_end = offset + len(line)
        offset = line_end
        if line.strip():
            segments.append(
                ParsedSegment(
                    text=line,
                    char_start=line_start,
                    char_end=line_end,
                    source_locator_json={"line_start": line_number, "line_end": line_number},
                    metadata_json={},
                )
            )

    if text and not segments:
        segments.append(
            ParsedSegment(
                text=text,
                char_start=0,
                char_end=len(text),
                source_locator_json={"line_start": 1, "line_end": 1},
                metadata_json={},
            )
        )

    return segments


def extract_markdown_heading_path(text: str) -> list[str]:
    path: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped.startswith("#"):
            continue

        level = len(stripped) - len(stripped.lstrip("#"))
        if level == 0 or level > 6 or not stripped[level:].startswith(" "):
            continue

        heading = stripped[level:].strip()
        path = path[: level - 1]
        path.append(heading)
        break

    return path


def build_structured_text(
    items: list[tuple[str, dict[str, Any], dict[str, Any]]],
) -> tuple[str, list[ParsedSegment]]:
    text_parts: list[str] = []
    segments: list[ParsedSegment] = []
    offset = 0
    for item_text, locator, metadata in items:
        if text_parts:
            text_parts.append("\n\n")
            offset += 2

        char_start = offset
        text_parts.append(item_text)
        offset += len(item_text)
        segments.append(
            ParsedSegment(
                text=item_text,
                char_start=char_start,
                char_end=offset,
                source_locator_json=locator,
                metadata_json=metadata,
            )
        )

    return "".join(text_parts), segments


def combine_source_locators(segments: list[ParsedSegment]) -> dict[str, Any]:
    if not segments:
        return {}
    if len(segments) == 1:
        return dict(segments[0].source_locator_json)

    first_locator = dict(segments[0].source_locator_json)
    first_locator["locators"] = [segment.source_locator_json for segment in segments]
    return first_locator


def combine_metadata(
    document_metadata: dict[str, Any],
    segments: list[ParsedSegment],
) -> dict[str, Any]:
    metadata = dict(document_metadata)
    for segment in segments:
        for key, value in segment.metadata_json.items():
            if value not in (None, [], {}, ""):
                metadata[key] = value
    return metadata


def package_version(package_name: str) -> str:
    try:
        return version(package_name)
    except PackageNotFoundError:
        return "unknown"
